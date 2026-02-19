#!/usr/bin/env python3
"""
Xcapit Privacy - Presentation Generator

Genera una presentación comercial profesional en PowerPoint.

Requisitos:
    pip install python-pptx
"""

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


# Brand colors
BRAND_PRIMARY = RGBColor(79, 70, 229)    # Indigo
BRAND_SECONDARY = RGBColor(16, 185, 129) # Green
NAVY = RGBColor(15, 23, 42)              # Navy
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(241, 245, 249)
RED = RGBColor(239, 68, 68)
DARK_GREEN = RGBColor(6, 95, 70)
DARK_RED = RGBColor(127, 29, 29)


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = NAVY
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = BRAND_SECONDARY
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, accent_color=None):
    """Add a content slide with bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = accent_color or BRAND_PRIMARY
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Bullet points
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = "• " + bullet
        p.font.size = Pt(22)
        p.font.color.rgb = NAVY
        p.space_before = Pt(12)
        p.space_after = Pt(12)

    return slide


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Add a two-column comparison slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column header
    left_header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.3), Inches(0.6)
    )
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = RED
    left_header.line.fill.background()

    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.55), Inches(4.3), Inches(0.5))
    tf = left_title_box.text_frame
    tf.paragraphs[0].text = left_title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Left bullets
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(4.3), Inches(4))
    tf = left_content.text_frame
    for i, bullet in enumerate(left_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "✗ " + bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_RED
        p.space_after = Pt(8)

    # Right column header
    right_header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(1.5), Inches(4.3), Inches(0.6)
    )
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = BRAND_SECONDARY
    right_header.line.fill.background()

    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.55), Inches(4.3), Inches(0.5))
    tf = right_title_box.text_frame
    tf.paragraphs[0].text = right_title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Right bullets
    right_content = slide.shapes.add_textbox(Inches(5.2), Inches(2.3), Inches(4.3), Inches(4))
    tf = right_content.text_frame
    for i, bullet in enumerate(right_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "✓ " + bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GREEN
        p.space_after = Pt(8)

    return slide


def add_stats_slide(prs, title, stats):
    """Add a statistics slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # Stats boxes
    num_stats = len(stats)
    box_width = 2.8
    total_width = num_stats * box_width + (num_stats - 1) * 0.3
    start_x = (10 - total_width) / 2

    for i, (value, label) in enumerate(stats):
        x = start_x + i * (box_width + 0.3)

        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2), Inches(box_width), Inches(2.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BRAND_PRIMARY

        # Value
        val_box = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(box_width), Inches(1))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = BRAND_PRIMARY
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(3.3), Inches(box_width - 0.2), Inches(1))
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_closing_slide(prs, title, subtitle, contact):
    """Add a closing CTA slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Accent line
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(4.8), prs.slide_width, Inches(0.1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND_SECONDARY
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(8), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = BRAND_SECONDARY
    p.alignment = PP_ALIGN.CENTER

    # Contact
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(0.5))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = contact
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide


def generate_presentation():
    """Generate the complete sales presentation."""
    print("🎨 Generando presentacion comercial...")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Xcapit Privacy",
        "Colaboracion de datos sin compartir datos"
    )
    print("  ✓ Slide 1: Portada")

    # Slide 2: The Problem
    add_content_slide(
        prs,
        "El Problema",
        [
            "Las empresas tienen datos valiosos pero no pueden compartirlos",
            "Regulaciones (GDPR, HIPAA) impiden compartir datos sensibles",
            "Los modelos de ML necesitan datos de multiples fuentes",
            "La falta de datos limita la innovacion y precision",
            "Los data breaches cuestan millones en multas"
        ],
        accent_color=RED
    )
    print("  ✓ Slide 2: El Problema")

    # Slide 3: Traditional vs Privacy-Preserving
    add_two_column_slide(
        prs,
        "Compartir Datos: Tradicional vs Privado",
        "Metodo Tradicional",
        [
            "Datos viajan en texto plano",
            "Servidor central ve todo",
            "Riesgo de data breach",
            "Problemas de compliance",
            "Confianza en terceros"
        ],
        "Xcapit Privacy",
        [
            "Datos siempre cifrados",
            "Servidor NO puede descifrar",
            "Imposible filtrar datos",
            "Compliance garantizado",
            "Zero-trust architecture"
        ]
    )
    print("  ✓ Slide 3: Comparacion")

    # Slide 4: How FHE Works
    add_content_slide(
        prs,
        "Encriptacion Homomorfica (FHE)",
        [
            "Los datos se cifran ANTES de salir de tu infraestructura",
            "El servidor procesa datos cifrados sin descifrarlos",
            "Operaciones matematicas funcionan sobre ciphertext",
            "Solo TU puedes descifrar el resultado final",
            "Esquema CKKS con 128 bits de seguridad"
        ],
        accent_color=BRAND_PRIMARY
    )
    print("  ✓ Slide 4: Como funciona FHE")

    # Slide 5: Use Cases
    add_content_slide(
        prs,
        "Casos de Uso",
        [
            "Fraude: Bancos colaboran sin compartir transacciones",
            "Credit Scoring: Datos de multiples bureaus",
            "Healthcare: Investigacion sin exponer pacientes",
            "Insurance: Modelos con datos de competidores",
            "Retail: Prediccion con datos de supply chain"
        ],
        accent_color=BRAND_SECONDARY
    )
    print("  ✓ Slide 5: Casos de Uso")

    # Slide 6: Stats
    add_stats_slide(
        prs,
        "Impacto Demostrado",
        [
            ("92%", "Precision del modelo"),
            ("0", "Datos expuestos"),
            ("100%", "Compliance garantizado"),
        ]
    )
    print("  ✓ Slide 6: Estadisticas")

    # Slide 7: How It Works (Flow)
    add_content_slide(
        prs,
        "Flujo de Trabajo",
        [
            "1. Crea un consorcio e invita participantes",
            "2. Cada empresa sube datos cifrados con su clave",
            "3. El modelo se entrena sobre datos cifrados",
            "4. Cada participante descarga sus predicciones",
            "5. Nadie ve los datos de nadie"
        ],
        accent_color=NAVY
    )
    print("  ✓ Slide 7: Flujo de Trabajo")

    # Slide 8: Security & Compliance
    add_content_slide(
        prs,
        "Seguridad y Compliance",
        [
            "ISO 27001 Certificado",
            "GDPR, HIPAA, SOC 2 Type II Compliant",
            "Encriptacion de 128 bits (nivel bancario)",
            "Auditoria completa de operaciones",
            "Zero-knowledge architecture"
        ],
        accent_color=DARK_GREEN
    )
    print("  ✓ Slide 8: Seguridad")

    # Slide 9: Closing
    add_closing_slide(
        prs,
        "Colabora sin comprometer",
        "El futuro de los datos es privado",
        "privacy.xcapit.com | contacto@xcapit.com"
    )
    print("  ✓ Slide 9: Cierre")

    # Save
    output_path = os.path.join(
        os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
        "Xcapit_Privacy_Comercial.pptx"
    )
    prs.save(output_path)

    print(f"\n✅ Presentacion guardada en:\n   {output_path}")
    return output_path


if __name__ == "__main__":
    generate_presentation()

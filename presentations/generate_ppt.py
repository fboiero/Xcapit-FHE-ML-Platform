#!/usr/bin/env python3
"""
Xcapit Privacy Platform - Commercial Presentation Generator
============================================================
Genera una presentacion PPT comercial con screenshots del demo
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import os

# Colores corporativos
XCAPIT_BLUE = RGBColor(79, 70, 229)  # Indigo-600
XCAPIT_GREEN = RGBColor(16, 185, 129)  # Emerald-500
DARK_TEXT = RGBColor(31, 41, 55)  # Gray-800
LIGHT_TEXT = RGBColor(107, 114, 128)  # Gray-500

BASE_DIR = Path(__file__).parent
SCREENSHOTS_DIR = BASE_DIR / "screenshots"

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background gradient effect (using shape)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = XCAPIT_BLUE
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(199, 210, 254)  # Indigo-200
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title, icon=""):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = XCAPIT_GREEN
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon} {title}" if icon else title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, image_path=None):
    """Add a content slide with bullets and optional image"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = XCAPIT_BLUE
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Content area
    if image_path and Path(image_path).exists():
        # Two column layout
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
        slide.shapes.add_picture(str(image_path), Inches(5.2), Inches(1.5), width=Inches(4.5))
    else:
        # Full width
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(12)

    return slide


def add_image_slide(prs, title, image_path, caption=""):
    """Add a full-image slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = XCAPIT_BLUE

    # Image
    if Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(1.2), width=Inches(9))
    else:
        # Placeholder
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(9), Inches(5))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(229, 231, 235)
        placeholder.line.color.rgb = RGBColor(156, 163, 175)

        text_box = slide.shapes.add_textbox(Inches(3), Inches(3.5), Inches(4), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = "[Screenshot placeholder]"
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = LIGHT_TEXT

    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = LIGHT_TEXT
        p.alignment = PP_ALIGN.CENTER

    return slide


def create_presentation():
    """Create the commercial presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ==========================================
    # SLIDE 1: Title
    # ==========================================
    add_title_slide(
        prs,
        "Xcapit Privacy Platform",
        "Machine Learning Colaborativo con Privacidad Garantizada"
    )

    # ==========================================
    # SLIDE 2: El Problema
    # ==========================================
    add_section_slide(prs, "El Problema", "")

    add_content_slide(prs, "El Dilema de los Datos", [
        "Las empresas tienen datos valiosos pero no pueden compartirlos",
        "Regulaciones estrictas (GDPR, LGPD, CCPA) limitan el uso de datos",
        "Los competidores no confian entre si para colaborar",
        "Los modelos de ML necesitan MAS datos para ser precisos",
        "El 73% de los datos empresariales nunca se utilizan por temas de privacidad"
    ])

    add_content_slide(prs, "Casos Reales de Friccion", [
        "Bancos que no pueden compartir datos de fraude entre si",
        "Hospitales con datos de pacientes aislados en silos",
        "Retailers que pierden insights por no colaborar",
        "Aseguradoras con modelos de riesgo suboptimos",
        "Fintechs limitadas por falta de historial crediticio"
    ])

    # ==========================================
    # SLIDE 5: La Solucion
    # ==========================================
    add_section_slide(prs, "La Solucion", "")

    add_content_slide(prs, "Cifrado Homomorfico (FHE)", [
        "Los datos NUNCA se descifran durante el procesamiento",
        "Computacion sobre datos cifrados - matematicamente seguro",
        "Ni siquiera el servidor puede ver los datos originales",
        "Cumple con todas las regulaciones de privacidad",
        "Auditado y verificable criptograficamente"
    ])

    add_content_slide(prs, "Aprendizaje Federado", [
        "Cada empresa mantiene sus datos en su infraestructura",
        "Solo se comparten gradientes/parametros cifrados",
        "El modelo aprende de todos sin ver datos individuales",
        "Escalable a multiples participantes",
        "Tolerante a fallas y desconexiones"
    ])

    # ==========================================
    # SLIDE 8: Demo - El Flujo
    # ==========================================
    add_section_slide(prs, "Demo en Vivo", "")

    add_image_slide(
        prs,
        "Paso 1: Datos Sensibles de Cada Cliente",
        SCREENSHOTS_DIR / "01_raw_data.png",
        "FinBank y RetailCorp tienen datos de clientes que NO pueden compartir directamente"
    )

    add_image_slide(
        prs,
        "Paso 2: Cifrado Homomorfico (FHE)",
        SCREENSHOTS_DIR / "02_encryption.png",
        "Los datos se cifran localmente - el servidor NUNCA ve los valores reales"
    )

    add_image_slide(
        prs,
        "Paso 3: Entrenamiento Federado",
        SCREENSHOTS_DIR / "03_training.png",
        "Random Forest entrenado sobre datos cifrados de ambos participantes"
    )

    add_image_slide(
        prs,
        "Paso 4: Resultados del Modelo",
        SCREENSHOTS_DIR / "04_results.png",
        "92% de accuracy - mejor que cualquier modelo individual"
    )

    add_image_slide(
        prs,
        "Paso 5: Prediccion en Tiempo Real",
        SCREENSHOTS_DIR / "05_prediction.png",
        "Nuevos clientes evaluados sin exponer datos historicos"
    )

    # ==========================================
    # SLIDE 14: Beneficios
    # ==========================================
    add_section_slide(prs, "Beneficios", "")

    add_content_slide(prs, "Para el Negocio", [
        "Modelos 40-60% mas precisos al combinar datos de multiples fuentes",
        "Nuevas oportunidades de colaboracion con competidores",
        "Monetizacion de datos sin riesgo de exposicion",
        "Ventaja competitiva en mercados regulados",
        "Reduccion de fraude y riesgo crediticio"
    ])

    add_content_slide(prs, "Para Compliance", [
        "100% compatible con GDPR, LGPD, CCPA",
        "Los datos nunca salen del control de cada empresa",
        "Audit trail completo y verificable",
        "Certificaciones de seguridad disponibles",
        "Demostracion matematica de privacidad"
    ])

    add_content_slide(prs, "Para IT/Seguridad", [
        "Integracion simple via API REST",
        "Sin cambios en infraestructura existente",
        "Cifrado end-to-end con claves propias",
        "Deployment on-premise o cloud",
        "SLA de 99.9% de disponibilidad"
    ])

    # ==========================================
    # SLIDE 18: Casos de Uso
    # ==========================================
    add_section_slide(prs, "Casos de Uso", "")

    add_content_slide(prs, "Servicios Financieros", [
        "Scoring crediticio colaborativo entre bancos",
        "Deteccion de fraude compartiendo patrones (no datos)",
        "Prevencion de lavado de dinero (AML) en consorcio",
        "Evaluacion de riesgo para seguros",
        "KYC colaborativo sin duplicar esfuerzos"
    ])

    add_content_slide(prs, "Salud y Life Sciences", [
        "Investigacion clinica multi-hospital",
        "Modelos de diagnostico con datos de multiples paises",
        "Farmacovigilancia colaborativa",
        "Prediccion de outcomes sin compartir PHI",
        "Genomica preservando privacidad del paciente"
    ])

    add_content_slide(prs, "Retail y CPG", [
        "Modelos de demanda con datos de multiples retailers",
        "Optimizacion de supply chain colaborativa",
        "Personalizacion sin compartir datos de clientes",
        "Deteccion de tendencias de mercado",
        "Pricing dinamico preservando secretos comerciales"
    ])

    # ==========================================
    # SLIDE 22: Pricing
    # ==========================================
    add_section_slide(prs, "Modelo Comercial", "")

    add_content_slide(prs, "Opciones de Licenciamiento", [
        "TIER 1 - Starter: $5K/mes - Hasta 3 participantes, 1M registros",
        "TIER 2 - Professional: $15K/mes - Hasta 10 participantes, 10M registros",
        "TIER 3 - Enterprise: $40K/mes - Ilimitado, SLA premium, soporte dedicado",
        "TIER 4 - Custom: Precio a medida para grandes consorcios",
        "POC gratuito de 30 dias con datos reales"
    ])

    # ==========================================
    # SLIDE 24: Proximos Pasos
    # ==========================================
    add_section_slide(prs, "Proximos Pasos", "")

    add_content_slide(prs, "Como Empezar", [
        "1. Workshop de descubrimiento (2 horas) - GRATIS",
        "2. POC con datos reales anonimizados (30 dias)",
        "3. Validacion tecnica y de compliance",
        "4. Piloto en produccion (3 meses)",
        "5. Rollout y escalamiento"
    ])

    # ==========================================
    # SLIDE 26: Contacto
    # ==========================================
    slide = add_title_slide(
        prs,
        "Gracias",
        "¿Preguntas?"
    )

    # Add contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "demo.xcapit.com | info@xcapit.com"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(199, 210, 254)
    p.alignment = PP_ALIGN.CENTER

    # Save
    output_path = BASE_DIR / "Xcapit_Privacy_Platform_Comercial.pptx"
    prs.save(str(output_path))
    print(f"Presentacion guardada en: {output_path}")

    return output_path


if __name__ == "__main__":
    create_presentation()
